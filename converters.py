"""
MasterConvert conversion engine.
Real conversions using LibreOffice (rendering-accurate formats) and
structural rebuilding (python-docx/pptx/openpyxl) for formats that have
no direct renderer path between them.
"""
import os
import subprocess
import tempfile
import shutil
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor as DocxRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt
from pptx.dml.color import RGBColor as PptxRGBColor
import openpyxl
from openpyxl.utils import get_column_letter
import pypdf


class ConversionError(Exception):
    pass


# ---------------------------------------------------------- PPTX templates
PPTX_TEMPLATES = {
    "minimal": {
        "bg": None, "title_fill": None,
        "title_font": "Calibri Light", "title_size": Pt(36), "title_bold": False,
        "title_color": PptxRGBColor(0x33, 0x33, 0x33),
        "body_font": "Calibri Light", "body_size": Pt(20),
        "body_color": PptxRGBColor(0x55, 0x55, 0x55),
    },
    "academic": {
        "bg": None, "title_fill": None,
        "title_font": "Georgia", "title_size": Pt(32), "title_bold": True,
        "title_color": PptxRGBColor(0x1F, 0x3A, 0x5F),
        "body_font": "Georgia", "body_size": Pt(18),
        "body_color": PptxRGBColor(0x22, 0x22, 0x22),
    },
    "bold": {
        "bg": PptxRGBColor(0x0A, 0x0A, 0x0A), "title_fill": None,
        "title_font": "Arial", "title_size": Pt(44), "title_bold": True,
        "title_color": PptxRGBColor(0xFF, 0xFF, 0xFF),
        "body_font": "Arial", "body_size": Pt(20),
        "body_color": PptxRGBColor(0xE0, 0xE0, 0xE0),
    },
    "classic": {
        "bg": None, "title_fill": PptxRGBColor(0x1F, 0x38, 0x64),
        "title_font": "Calibri", "title_size": Pt(30), "title_bold": True,
        "title_color": PptxRGBColor(0xFF, 0xFF, 0xFF),
        "body_font": "Calibri", "body_size": Pt(18),
        "body_color": PptxRGBColor(0x22, 0x22, 0x22),
    },
}


def _style_pptx_slide(slide, style):
    if style["bg"] is not None:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = style["bg"]
    title_shape = slide.shapes.title
    if style["title_fill"] is not None:
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = style["title_fill"]
    for para in title_shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = style["title_font"]
            run.font.size = style["title_size"]
            run.font.bold = style["title_bold"]
            run.font.color.rgb = style["title_color"]


def _style_pptx_body_paragraph(p, style):
    p.font.name = style["body_font"]
    p.font.size = style["body_size"]
    p.font.color.rgb = style["body_color"]


# ---------------------------------------------------------- DOCX styles
DOCX_STYLES = {
    "clean": {
        "heading_font": "Calibri", "heading_color": DocxRGBColor(0x22, 0x22, 0x22),
        "body_font": "Calibri", "body_size": DocxPt(11), "uppercase_headings": False,
        "page_numbers": False, "tight_spacing": False,
    },
    "academic": {
        "heading_font": "Cambria", "heading_color": DocxRGBColor(0x1F, 0x3A, 0x5F),
        "body_font": "Cambria", "body_size": DocxPt(12), "uppercase_headings": False,
        "page_numbers": True, "tight_spacing": False,
    },
    "report": {
        "heading_font": "Calibri", "heading_color": DocxRGBColor(0x7A, 0x1F, 0x2B),
        "body_font": "Calibri", "body_size": DocxPt(11), "uppercase_headings": True,
        "page_numbers": True, "tight_spacing": False,
    },
    "compact": {
        "heading_font": "Calibri", "heading_color": DocxRGBColor(0x22, 0x22, 0x22),
        "body_font": "Calibri", "body_size": DocxPt(9), "uppercase_headings": False,
        "page_numbers": False, "tight_spacing": True,
    },
}


def _style_docx_headings(doc, style):
    """Applies to every heading paragraph already in the document (headings
    are added via doc.add_heading before this runs)."""
    for para in doc.paragraphs:
        if not para.style.name.lower().startswith("heading") and para.style.name != "Title":
            continue
        if style["uppercase_headings"]:
            for run in para.runs:
                run.text = run.text.upper()
        for run in para.runs:
            run.font.name = style["heading_font"]
            run.font.color.rgb = style["heading_color"]


def _style_docx_body(doc, style):
    for para in doc.paragraphs:
        is_heading = para.style.name.lower().startswith("heading") or para.style.name == "Title"
        if is_heading:
            continue
        for run in para.runs:
            run.font.name = style["body_font"]
            run.font.size = style["body_size"]
        if style.get("tight_spacing"):
            para.paragraph_format.space_before = DocxPt(0)
            para.paragraph_format.space_after = DocxPt(2)


def _add_docx_page_number_footer(doc):
    section = doc.sections[0]
    footer = section.footer
    para = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
    para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = para.add_run()
    fld_begin = OxmlElement("w:fldChar")
    fld_begin.set(qn("w:fldCharType"), "begin")
    instr = OxmlElement("w:instrText")
    instr.set(qn("xml:space"), "preserve")
    instr.text = "PAGE"
    fld_end = OxmlElement("w:fldChar")
    fld_end.set(qn("w:fldCharType"), "end")
    run._r.append(fld_begin)
    run._r.append(instr)
    run._r.append(fld_end)


def apply_docx_style(doc, style_name):
    style = DOCX_STYLES.get(style_name, DOCX_STYLES["clean"])
    _style_docx_headings(doc, style)
    _style_docx_body(doc, style)
    if style["page_numbers"]:
        _add_docx_page_number_footer(doc)


def _soffice_convert(src_path, target_format, out_dir):
    """Use headless LibreOffice for true rendering-based conversions
    (currently: docx->pdf). Raises ConversionError on failure."""
    result = subprocess.run(
        ["soffice", "--headless", "--nologo", "--nofirststartwizard",
         "--convert-to", target_format, "--outdir", out_dir, src_path],
        capture_output=True, text=True, timeout=60
    )
    base = os.path.splitext(os.path.basename(src_path))[0]
    out_path = os.path.join(out_dir, f"{base}.{target_format}")
    if not os.path.exists(out_path):
        raise ConversionError(f"LibreOffice conversion failed: {result.stderr or result.stdout}")
    return out_path


# ---------------------------------------------------------------- DOCX -> PDF
def docx_to_pdf(src_path, out_dir, style=None):
    return _soffice_convert(src_path, "pdf", out_dir)


# --------------------------------------------------------------- DOCX -> PPTX
def docx_to_pptx(src_path, out_dir, style="minimal"):
    template_style = PPTX_TEMPLATES.get(style, PPTX_TEMPLATES["minimal"])
    doc = Document(src_path)
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)
    title_layout = prs.slide_layouts[1]  # title + content

    slide = None
    body_tf = None

    def new_slide(title_text):
        s = prs.slides.add_slide(title_layout)
        s.shapes.title.text = title_text or "Untitled"
        _style_pptx_slide(s, template_style)
        tf = s.placeholders[1].text_frame
        tf.clear()
        return s, tf

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        para_style_name = (para.style.name or "").lower()
        if "heading" in para_style_name or para_style_name == "title":
            slide, body_tf = new_slide(text)
        else:
            if slide is None:
                slide, body_tf = new_slide("Overview")
            if body_tf.paragraphs[0].text == "" and len(body_tf.paragraphs) == 1:
                p = body_tf.paragraphs[0]
            else:
                p = body_tf.add_paragraph()
            p.text = text
            p.level = 0
            _style_pptx_body_paragraph(p, template_style)

    if slide is None:
        new_slide("Untitled Document")

    out_path = os.path.join(out_dir, "converted.pptx")
    prs.save(out_path)
    return out_path


# --------------------------------------------------------------- PPTX -> DOCX
def pptx_to_docx(src_path, out_dir, style="clean"):
    prs = Presentation(src_path)
    doc = Document()
    doc.add_heading("Slide Handout", 0)

    for i, slide in enumerate(prs.slides, 1):
        title = None
        body_texts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if shape == slide.shapes.title:
                title = text
            else:
                body_texts.append(text)

        doc.add_heading(title or f"Slide {i}", level=1)
        for t in body_texts:
            for line in t.split("\n"):
                if line.strip():
                    doc.add_paragraph(line.strip(), style="List Bullet")
        if i < len(prs.slides):
            doc.add_page_break()

    apply_docx_style(doc, style)
    out_path = os.path.join(out_dir, "converted.docx")
    doc.save(out_path)
    return out_path


# ---------------------------------------------------------------- PDF -> DOCX
def pdf_to_docx(src_path, out_dir, style="clean"):
    reader = pypdf.PdfReader(src_path)
    doc = Document()
    doc.add_heading("Converted from PDF", 0)

    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        text = text.strip()
        if len(reader.pages) > 1:
            doc.add_heading(f"Page {i}", level=2)
        if text:
            for line in text.split("\n"):
                line = line.strip()
                if line:
                    doc.add_paragraph(line)
        else:
            doc.add_paragraph("[No extractable text on this page — likely a scanned image.]")

    apply_docx_style(doc, style)
    out_path = os.path.join(out_dir, "converted.docx")
    doc.save(out_path)
    return out_path


# ---------------------------------------------------------------- PDF -> PPTX
def pdf_to_pptx(src_path, out_dir, style=None):
    page_prefix = os.path.join(out_dir, "page")
    subprocess.run(
        ["pdftoppm", "-png", "-r", "150", src_path, page_prefix],
        capture_output=True, text=True, timeout=60, check=True
    )
    page_images = sorted(
        f for f in os.listdir(out_dir) if f.startswith("page") and f.endswith(".png")
    )
    if not page_images:
        raise ConversionError("Could not rasterize PDF pages")

    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)
    blank_layout = prs.slide_layouts[6]

    for img_name in page_images:
        img_path = os.path.join(out_dir, img_name)
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

        ocr_text = ""
        try:
            ocr = subprocess.run(
                ["tesseract", img_path, "-", "--psm", "6"],
                capture_output=True, text=True, timeout=30
            )
            ocr_text = ocr.stdout.strip()
        except Exception:
            pass
        notes = slide.notes_slide
        notes.notes_text_frame.text = (
            ocr_text if ocr_text else "(No text detected by OCR on this page.)"
        )

    out_path = os.path.join(out_dir, "converted.pptx")
    prs.save(out_path)
    return out_path


# --------------------------------------------------------------- XLSX -> DOCX
def xlsx_to_docx(src_path, out_dir, style="clean"):
    wb = openpyxl.load_workbook(src_path, data_only=True)
    doc = Document()
    doc.add_heading("Converted from Excel", 0)

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        doc.add_heading(sheet_name, level=1)
        rows = list(ws.iter_rows(values_only=True))
        rows = [r for r in rows if any(c is not None and str(c).strip() for c in r)]
        if not rows:
            doc.add_paragraph("(Empty sheet)")
            continue
        n_cols = max(len(r) for r in rows)
        table = doc.add_table(rows=len(rows), cols=n_cols)
        table.style = "Light Grid Accent 1"
        for r_idx, row in enumerate(rows):
            for c_idx in range(n_cols):
                val = row[c_idx] if c_idx < len(row) and row[c_idx] is not None else ""
                table.cell(r_idx, c_idx).text = str(val)

    apply_docx_style(doc, style)
    out_path = os.path.join(out_dir, "converted.docx")
    doc.save(out_path)
    return out_path


# --------------------------------------------------------------- DOCX -> XLSX
def docx_to_xlsx(src_path, out_dir):
    doc = Document(src_path)
    wb = openpyxl.Workbook()
    wb.remove(wb.active)

    if doc.tables:
        for i, table in enumerate(doc.tables, 1):
            ws = wb.create_sheet(title=f"Table {i}"[:31])
            for r_idx, row in enumerate(table.rows, 1):
                for c_idx, cell in enumerate(row.cells, 1):
                    ws.cell(row=r_idx, column=c_idx, value=cell.text)
            for col in range(1, len(table.columns) + 1):
                ws.column_dimensions[get_column_letter(col)].width = 22
    else:
        ws = wb.create_sheet(title="Document Text")
        ws.column_dimensions["A"].width = 100
        r = 1
        for para in doc.paragraphs:
            if para.text.strip():
                ws.cell(row=r, column=1, value=para.text.strip())
                r += 1

    out_path = os.path.join(out_dir, "converted.xlsx")
    wb.save(out_path)
    return out_path


CONVERTERS = {
    ("docx", "pdf"): docx_to_pdf,
    ("docx", "pptx"): docx_to_pptx,
    ("pptx", "docx"): pptx_to_docx,
    ("pdf", "docx"): pdf_to_docx,
    ("pdf", "pptx"): pdf_to_pptx,
    ("xlsx", "docx"): xlsx_to_docx,
    ("docx", "xlsx"): docx_to_xlsx,
}


def convert(src_path, from_fmt, to_fmt, out_dir, style=None):
    key = (from_fmt.lower(), to_fmt.lower())
    if key not in CONVERTERS:
        raise ConversionError(f"Unsupported conversion: {from_fmt} -> {to_fmt}")
    if style:
        return CONVERTERS[key](src_path, out_dir, style=style)
    return CONVERTERS[key](src_path, out_dir)


# ------------------------------------------------------- raw text -> PPTX
def text_to_pptx(raw_text, out_dir):
    """Turn pasted plain text into a slide deck. Blocks separated by a
    blank line become slides; each block's first line is the slide title,
    remaining lines become bullets. No AI involved — pure structural rules,
    same spirit as docx_to_pptx but for text with no formatting to read."""
    blocks = [b.strip() for b in raw_text.strip().split("\n\n") if b.strip()]
    if not blocks:
        raise ConversionError("No text provided")

    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)
    title_layout = prs.slide_layouts[1]

    for block in blocks:
        lines = [l.strip() for l in block.split("\n") if l.strip()]
        if not lines:
            continue
        title, body_lines = lines[0], lines[1:]

        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = title[:120]
        tf = slide.placeholders[1].text_frame
        tf.clear()
        if body_lines:
            for i, line in enumerate(body_lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = line

    if not prs.slides:
        raise ConversionError("No usable text found")

    out_path = os.path.join(out_dir, "Presentation.pptx")
    prs.save(out_path)
    return out_path


# --------------------------------------------------------- extract raw text
def extract_text(src_path, ext):
    """Pull plain text out of an uploaded file for use in the AI tools.
    Supports the same four formats the converter already handles, plus .txt."""
    ext = ext.lower()
    if ext == "txt":
        with open(src_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    if ext == "docx":
        doc = Document(src_path)
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                parts.append(" | ".join(c.text for c in row.cells))
        return "\n".join(parts)
    if ext == "pdf":
        reader = pypdf.PdfReader(src_path)
        return "\n".join((page.extract_text() or "") for page in reader.pages)
    if ext == "pptx":
        prs = Presentation(src_path)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    parts.append(shape.text_frame.text)
        return "\n".join(parts)
    if ext == "xlsx":
        wb = openpyxl.load_workbook(src_path, data_only=True)
        parts = []
        for sheet in wb.sheetnames:
            for row in wb[sheet].iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts)
    raise ConversionError(f"Can't extract text from .{ext} files")
